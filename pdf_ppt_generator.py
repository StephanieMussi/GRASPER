import os
import sys
import fitz  
import copy 
import math
import time
import torch
import argparse
import threading
import tkinter as tk
from openai import OpenAI
from llama_cpp import Llama
from pptx import Presentation
from tkinter import filedialog
from transformers import BartTokenizer, BartForConditionalGeneration
from transformers import T5Tokenizer, T5ForConditionalGeneration
from transformers import LlamaForCausalLM, LlamaTokenizer
from transformers import AutoModelForCausalLM, AutoTokenizer

"""
Validate the arguments specified by user. 
If invalid arguments are given, terminate the program.
"""
def validate_arguments(args):
    if args.summarizer != 'bart' and args.summarizer != 't5' and args.summarizer != 'llama2' and args.summarizer != 'mistral':
        sys.exit("Error: No valid summarizer specified. Please use '--summarizer=bart' or '--summarizer=t5' or '--summarizer=llama2' or '--summarizer=mistral'.")
    
    if args.converter == 'gpt3':
        if args.openaikey is None:
            sys.exit("Error: No OpenAI API key provided. Please use '--openaikey=<your-openai-key>'.")
    elif args.converter != 'llama2' and  args.converter != 'mistral':
        sys.exit("Error: No valid converter specified. Please use '--converter=gpt3' or '--converter=llama2' or '--converter=mistral'.")


"""
Get pdf file input in pop-up file manager window.
If no file is selected, terminate the program.
"""
def ask_for_file():
    print("Please upload paper in pdf format:")
    root = tk.Tk()
    root.withdraw() 
    file_path = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
    if file_path:
        print(f"File selected: {file_path}")
    else:
        sys.exit("No file selected. Exiting...")
    return file_path


"""
Get paper title from pdf file path.
It is assumed that name of the pdf file is the paper title.
"""
def get_paper_title(pdf_path):
    filename_with_extension = os.path.basename(pdf_path)
    paper_title, extension = os.path.splitext(filename_with_extension)
    return paper_title


"""
Extract the text from pdf file with PyMuPDF
"""
def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)  
    all_text = ""  
    for page in doc:
        all_text += page.get_text() 
    doc.close() 
    return all_text


"""
Split text into chunks with specified number of tokens.
"""
def chunk_text_into_tokens(text, tokenizer, max_tokens=1024):
    tokens = tokenizer.tokenize(text)
    token_chunks = [tokens[i:i + max_tokens] for i in range(0, len(tokens), max_tokens)]
    count = len(token_chunks)
    return token_chunks, count


"""
With specified maximum summary length, calculate the min and max length for each chunk summary.
"""
def get_max_min_sum_length(num_chunk, max_text_length=3996):
    max_sum_length = math.floor(max_text_length / num_chunk)
    min_sum_length = math.floor(0.5 * max_sum_length)
    return max_sum_length, min_sum_length


"""
Join chunk summaries list into one single string
"""
def concatenate_summaries(summaries):
    return " ".join(summaries)


"""
Summarize text with BART model
"""
def summarize_chunks_bart(token_chunks, max_sum_length, min_sum_length, tokenizer, model):
    summaries = []
    for chunk in token_chunks:
        chunk_text = tokenizer.convert_tokens_to_string(chunk)
        inputs = tokenizer(chunk_text, return_tensors='pt', padding=True, truncation=True, max_length=1024)
        summary_ids = model.generate(**inputs, max_length=max_sum_length, min_length=min_sum_length, num_beams=4, early_stopping=True)
        summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
        summaries.append(summary)
    return summaries

def summarize_text_bart(long_text, max_text_length=3996):
    tokenizer_bart = BartTokenizer.from_pretrained("facebook/bart-large-cnn")
    model_bart = BartForConditionalGeneration.from_pretrained("facebook/bart-large-cnn")
    text_chunks, num_chunk = chunk_text_into_tokens(long_text, tokenizer_bart)
    max_sum_length, min_sum_length = get_max_min_sum_length(num_chunk)
    chunk_summaries = summarize_chunks_bart(text_chunks, max_sum_length, min_sum_length, tokenizer_bart, model_bart)
    final_summary = concatenate_summaries(chunk_summaries)
    return final_summary


"""
Summarize text with T5 model
"""
def summarize_chunks_T5(chunks, max_sum_length, min_sum_length, tokenizer, model):
    summaries = []
    for chunk in chunks:
        chunk_text = tokenizer.convert_tokens_to_string(chunk)
        inputs = tokenizer("summarize: "+chunk_text, return_tensors='pt', padding=True, truncation=True, max_length=1024)
        summary_ids = model.generate(**inputs, max_length=max_sum_length, min_length=min_sum_length, num_beams=4, early_stopping=True)
        summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
        summaries.append(summary)
    return summaries

def summarize_text_T5(long_text, max_text_length=3996):
    tokenizer_T5= T5Tokenizer.from_pretrained("t5-base", legacy=False) 
    model_T5 = T5ForConditionalGeneration.from_pretrained("t5-base")
    text_chunks, num_chunk = chunk_text_into_tokens(long_text, tokenizer_T5)
    max_sum_length, min_sum_length = get_max_min_sum_length(num_chunk)
    chunk_summaries_T5 = summarize_chunks_T5(text_chunks, max_sum_length, min_sum_length, tokenizer_T5, model_T5)
    final_summary_T5 = concatenate_summaries(chunk_summaries_T5)
    return final_summary_T5


"""
Summarize text with LlaMa2 model
"""
def summarize_chunks_llama2(token_chunks, max_sum_length, min_sum_length, tokenizer):
    summaries = []
    for chunk in token_chunks:
        chunk_text = tokenizer.convert_tokens_to_string(chunk)
        llm = Llama(
            model_path="./models/llama-2-7b-chat.Q4_K_M.gguf",
            chat_format="llama-2", 
            n_ctx=4096
        )
        completion_llama2 = llm.create_chat_completion(
            messages=[
                {"role": "system", "content": chunk_text},
                {
                    "role": "user",
                    "content": "Please summarize the text between {} and {} tokens.".format(min_sum_length, max_sum_length)
                }
            ],
            max_tokens=max_sum_length
        )
        message_llama2 = completion_llama2['choices'][0]['message']['content']
        summaries.append(message_llama2)
    return summaries

def summarize_text_llama2(long_text):
    tokenizer_llama2 = LlamaTokenizer.from_pretrained("meta-llama/Llama-2-7b-hf")
    text_chunks, num_chunk = chunk_text_into_tokens(long_text, tokenizer_llama2, 3996)
    max_sum_length, min_sum_length = get_max_min_sum_length(num_chunk)
    chunk_summaries_llama2 = summarize_chunks_llama2(text_chunks, max_sum_length, min_sum_length, tokenizer_llama2)
    final_summary_llama2 = concatenate_summaries(chunk_summaries_llama2)
    return final_summary_llama2


"""
Summarize text with Mistral model
"""
def summarize_chunks_mistral(token_chunks, max_sum_length, min_sum_length, tokenizer):
    summaries = []
    for chunk in token_chunks:
        chunk_text = tokenizer.convert_tokens_to_string(chunk)
        mistral = Llama(
            model_path="./models/mistral-7b-instruct-v0.2.Q4_K_M.gguf",
            chat_format="llama-2", 
            n_ctx=4096
        )
        completion_mistral = mistral.create_chat_completion(
            messages = [
                {"role": "system", "content": chunk_text},
                {
                    "role": "user",
                    "content": "Please summarize the text between {} and {} tokens.".format(min_sum_length, max_sum_length)
                }
            ],
            max_tokens=max_sum_length
        )
        message_mistral = completion_mistral['choices'][0]['message']['content']
        summaries.append(message_mistral)
    return summaries

def summarize_text_mistral(long_text):
    tokenizer_mistral = AutoTokenizer.from_pretrained("mistralai/Mistral-7B-v0.1")
    text_chunks, num_chunk = chunk_text_into_tokens(long_text, tokenizer_mistral, 3996)
    max_sum_length, min_sum_length = get_max_min_sum_length(num_chunk)
    chunk_summaries_mistral = summarize_chunks_mistral(text_chunks, max_sum_length, min_sum_length, tokenizer_mistral)
    final_summary_mistral = concatenate_summaries(chunk_summaries_mistral)
    return final_summary_mistral


"""
Call corresponding summarizer function based on argument.
"""
def summarize_text(long_text, summarizer_type):
    original_stderr = sys.stderr
    sys.stderr = open(os.devnull, 'w')
    try:
        if summarizer_type == 'bart':
            return summarize_text_bart(long_text)
        elif summarizer_type == 't5':
            return summarize_text_T5(long_text)
        elif summarizer_type == 'llama2':
            return summarize_text_llama2(long_text)
        elif summarizer_type == 'mistral':
            return summarize_text_mistral(long_text)
    finally:
        sys.stderr.close()
        sys.stderr = original_stderr


"""
Convert summary to slides content with GPT3
"""
def convert_content_gpt3(final_summary, openaikey):
    client = OpenAI(api_key = openaikey)
    completion = client.chat.completions.create(
    model="gpt-3.5-turbo",
    messages=[
        {"role": "system", "content": final_summary},
        {"role": "user", "content": "Please convert the text into a presentation slides. Give the title and actual detailed bullet point content for each slide, in the format of \"Slide 1\nTitle: (title)\nContent: (content)\". Do not add other information."}
    ]
    )
    message = completion.choices[0].message.content
    return message


"""
Convert summary to slides content with LlaMa2
"""
def convert_content_llama2(final_summary):
    llm = Llama(
        model_path="./models/llama-2-7b-chat.Q4_K_M.gguf",
        chat_format="llama-2", 
        n_ctx=4096
    )
    completion_llama2 = llm.create_chat_completion(
        messages=[
            {"role": "system", "content": final_summary},
            {
                "role": "user",
                "content": "Please convert the text into presentation slides. Give the title and actual detailed bullet point content for each slide, in the format of 'Slide 1\nTitle: (title)\nContent: (content)'. Do not add other information."
            }
        ]
    )
    message_llama2 = completion_llama2['choices'][0]['message']['content']
    return message_llama2


"""
Convert summary to slides content with Mistral
"""
def convert_content_mistral(final_summary):
    mistral = Llama(
        model_path="./models/mistral-7b-instruct-v0.2.Q4_K_M.gguf",
        chat_format="llama-2", 
        n_ctx=4096
    )
    completion_mistral = mistral.create_chat_completion(
        messages = [
            {"role": "system", "content": final_summary},
            {
                "role": "user",
                "content": "Please convert the text into a presentation slides. Give the title and actual detailed bullet point content for each slide, in the format of \"Slide 1\nTitle: (title)\nContent: (content)\". Do not add other information. "
            }
        ]
    )
    message_mistral = completion_mistral['choices'][0]['message']['content']
    return message_mistral


"""
Call corresponding converter function based on argument.
"""
def convert_content(final_summary, converter_type, openaikey=None):
    original_stderr = sys.stderr
    sys.stderr = open(os.devnull, 'w')
    try:
        if converter_type == 'gpt3':
            return convert_content_gpt3(final_summary, openaikey)
        elif converter_type == 'llama2':
            return convert_content_llama2(final_summary)
        elif converter_type == 'mistral':
            return convert_content_mistral(final_summary)
    finally:
        sys.stderr.close()
        sys.stderr = original_stderr


"""
Add title slide by replacing content in template title slide
"""
def add_title_slide(title, prs):
    first_slide = prs.slides[0]
    title_shape = first_slide.shapes[0]
    if title_shape.has_text_frame:
        p = title_shape.text_frame.paragraphs[0]  
        run = p.runs[0] if p.runs else p.add_run()  
        run.text = title 


"""
Add content slide by replacing content in template content slide
"""       
def add_content_slide(title, content, prs):
    first_slide = prs.slides[-1]
    title_shape = first_slide.shapes[0]
    if title_shape.has_text_frame:
        p = title_shape.text_frame.paragraphs[0]  
        run = p.runs[0] if p.runs else p.add_run()  
        run.text = title 
    content_shape = first_slide.shapes[1]
    if content_shape.has_text_frame:
        p = content_shape.text_frame.paragraphs[0]  
        run = p.runs[0] if p.runs else p.add_run()  
        run.text = content 


"""
Duplicate template content slides to fit the number of content slides in slides content text
"""
def duplicate_slide(presentation, slide_index):
    slide = presentation.slides[slide_index]
    layout = slide.slide_layout
    duplicated_slide = presentation.slides.add_slide(layout)
    for shape in slide.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)
        duplicated_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return presentation


"""
Convert slides content text into ppt file.
It is assumed that the first slide is title slide, and rest are content slides.
"""
def generate_slides(paper_title, message, template_path, save_path):
    slides = message.strip().split("Slide ")[1:]
    prs = Presentation(template_path)
    add_title_slide(paper_title, prs)
    for i, slide in enumerate(slides):
        parts = slide.split("\nTitle:")
        slide_number = parts[0].strip()  
        title_content = parts[1].split("\nContent:")
        title = title_content[0].strip()
        content = title_content[1].strip()
        if i > 0:
            prs = duplicate_slide(prs, -1)
        add_content_slide(title, content, prs)     
    prs.save(save_path)


"""
Display dynamic progress indicator in terminal.
This is to let users know the current progress of slides generation.
"""
def animate_extraction():
    global extracting
    i = 0
    while extracting:
        sys.stdout.write("\Extracting text" + "." * i + " " * (3 - i)) 
        sys.stdout.flush()
        time.sleep(0.5)  
        i = (i + 1) % 4 
    sys.stdout.write("\r" + " " * 20 + "\r")  
    sys.stdout.flush()

def animate_summarization():
    global summarizing
    i = 0
    while summarizing:
        sys.stdout.write("\rSummarizing text" + "." * i + " " * (3 - i)) 
        sys.stdout.flush()
        time.sleep(0.5)  
        i = (i + 1) % 4 
    sys.stdout.write("\r" + " " * 20 + "\r")  
    sys.stdout.flush()

def animate_conversion():
    global converting
    i = 0
    while converting:
        sys.stdout.write("\rConverting content" + "." * i + " " * (3 - i))  
        sys.stdout.flush()
        time.sleep(0.5) 
        i = (i + 1) % 4  
    sys.stdout.write("\r" + " " * 30 + "\r") 
    sys.stdout.flush()

def animate_generation():
    global generating
    i = 0
    while generating:
        sys.stdout.write("\rGenerating slides" + "." * i + " " * (3 - i)) 
        sys.stdout.flush()
        time.sleep(0.5)  
        i = (i + 1) % 4 
    sys.stdout.write("\r" + " " * 20 + "\r")  
    sys.stdout.flush()


"""
Main function
"""   
def main():
    """Parse arguments"""
    parser = argparse.ArgumentParser(description="Process some arguments.")
    parser.add_argument('--summarizer', type=str, help='Choose the model summarizer (BART or T5 or llama2 or mistral2)')
    parser.add_argument('--converter', type=str, help='Choose the model converter (gpt3 or llama2 or mistral2)')
    parser.add_argument('--openaikey', type=str, help='Provide your OpenAI API key to access GPT3 API')
    args = parser.parse_args()
    validate_arguments(args)

    """Get pdf file from upload"""
    pdf_path = ask_for_file()
    paper_title = get_paper_title(pdf_path)

    """Step1. Extract text from pdf"""
    global extracting
    extracting = True
    anim_thread = threading.Thread(target=animate_extraction)
    anim_thread.start()
    long_text = extract_text_from_pdf(pdf_path)
    extracting = False
    anim_thread.join() 

    """Step2. Summarize text to reduce length"""
    global summarizing
    summarizing = True
    anim_thread = threading.Thread(target=animate_summarization)
    anim_thread.start()
    final_summary = summarize_text(long_text, args.summarizer)
    summarizing = False
    anim_thread.join() 

    """Step3. Convert summary into slides content"""
    global converting
    converting = True
    conv_thread = threading.Thread(target=animate_conversion)
    conv_thread.start()
    message = convert_content(final_summary, args.converter, args.openaikey)
    converting = False
    conv_thread.join() 

    """Step4. Generate presentation slides from content"""
    template_path = './source/template.pptx'
    save_path = './output/'+args.summarizer+' '+args.converter+'/'+paper_title+'.pptx'
    global generating
    generating = True
    anim_thread = threading.Thread(target=animate_generation)
    anim_thread.start()
    generate_slides(paper_title, message, template_path, save_path)
    generating = False
    anim_thread.join() 
    print("Slides is saved to: " + os.path.abspath(save_path))


if __name__ == "__main__":
    main()